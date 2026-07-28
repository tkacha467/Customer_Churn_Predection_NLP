from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # 0 - Title (with subtitle)
    # 1 - Title and Content
    # 5 - Title Only
    
    # SLIDE 1: TITLE
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "ChurnLens: Customer Risk Intelligence"
    subtitle.text = "Predictive Analytics & NLP Review Integrity Dashboard\nGenerated automatically from the Data Engineering Pipeline"

    # SLIDE 2: OVERVIEW
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Project Overview & Architecture"
    tf = content.text_frame
    tf.text = "ChurnLens is an end-to-end Machine Learning pipeline that predicts customer churn and evaluates review integrity."
    p = tf.add_paragraph()
    p.text = "1. Data Engineering: 100k+ rows merged and cleaned."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Feature Engineering: RFM metrics (Recency, Frequency, Monetary) extracted."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Predictive Modeling: XGBoost model handles imbalanced churn data using SMOTE."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. NLP Integrity: VADER + TF-IDF Logistic Regression detects mismatching reviews."
    p.level = 1

    # SLIDE 3: XGBOOST CONFUSION MATRIX
    if os.path.exists('models/confusion_matrix.png'):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "XGBoost Churn Prediction Performance"
        
        # Add Image
        img_path = 'models/confusion_matrix.png'
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(2), width=Inches(7))

    # SLIDE 4: FEATURE IMPORTANCE
    if os.path.exists('models/feature_importance.png'):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Key Drivers of Customer Churn"
        
        img_path = 'models/feature_importance.png'
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

    # SLIDE 5: NLP INTEGRITY
    if os.path.exists('models/nlp_confusion_matrix.png'):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "NLP Review Integrity Detection"
        
        # Text explanation
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = "Model successfully identifies cases where the text sentiment contradicts the 5-star or 1-star rating."
        
        img_path = 'models/nlp_confusion_matrix.png'
        slide.shapes.add_picture(img_path, Inches(2), Inches(2.5), width=Inches(6))

    # SLIDE 6: OVERALL RISK
    if os.path.exists('models/risk_distribution.png'):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Overall Customer Risk Distribution"
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = "Final Ensembled Risk Score combining XGBoost Churn Probability and NLP Integrity anomalies."
        
        img_path = 'models/risk_distribution.png'
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(2.5), width=Inches(7))

    # Save
    prs.save('ChurnLens_Presentation.pptx')
    print("Presentation saved as ChurnLens_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
